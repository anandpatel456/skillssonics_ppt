"""
Diagram Generator
Creates simple architecture / flow diagrams on a slide.

Fixed:
  - Nodes were positioned without checking slide width → overflow
  - Text inside nodes had no colour set → defaulted to black on dark shapes
  - No font size set → inconsistent rendering
"""

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# How wide one node box is
NODE_W = 1.7
NODE_H = 0.7
GAP    = 0.25   # horizontal gap between nodes
START_X = 0.4
START_Y = 3.2   # vertical centre of the diagram row


class DiagramGenerator:

    def create_architecture(self, slide, nodes: list):
        """
        Draw a horizontal row of rounded-rectangle nodes with arrows.
        Automatically wraps to a second row if nodes would overflow the slide.
        """
        if not nodes:
            return

        colors = ["#1E3A5F", "#F5A623", "#27AE60", "#1E3A5F", "#F5A623", "#27AE60"]
        max_per_row = int((10.0 - START_X) // (NODE_W + GAP))   # ~5 nodes per row

        for i, node in enumerate(nodes):
            row = i // max_per_row
            col = i %  max_per_row

            x = START_X + col * (NODE_W + GAP)
            y = START_Y + row * (NODE_H + 0.5)

            color_hex = colors[i % len(colors)]
            rgb = tuple(int(color_hex.lstrip("#")[j:j+2], 16) for j in (0, 2, 4))

            # node box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(NODE_W), Inches(NODE_H),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*rgb)
            box.line.fill.background()

            # node text
            tb = slide.shapes.add_textbox(
                Inches(x + 0.05), Inches(y + 0.08),
                Inches(NODE_W - 0.1), Inches(NODE_H - 0.1),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.text = node
            p = tf.paragraphs[0]
            p.font.size      = Pt(11)
            p.font.bold      = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment      = PP_ALIGN.CENTER

            # arrow between nodes on the same row (not after last in row)
            is_last_in_row = (col == max_per_row - 1) or (i == len(nodes) - 1)
            if not is_last_in_row:
                ax = x + NODE_W + 0.03
                ay = y + NODE_H / 2 - 0.13
                arr = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(ax), Inches(ay),
                    Inches(GAP - 0.05), Inches(0.26),
                )
                arr.fill.solid()
                arr.fill.fore_color.rgb = RGBColor(160, 160, 160)
                arr.line.fill.background()