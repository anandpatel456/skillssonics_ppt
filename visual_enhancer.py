"""
Visual Enhancer — produces unique visuals for each slide.
Uses slide index + topic to ensure different images per slide.
"""

import os, hashlib, urllib.request, urllib.parse, urllib.error, json, random
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

CACHE_DIR = "/tmp/ss_img"
os.makedirs(CACHE_DIR, exist_ok=True)

# ── Colours ───────────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
C_GOLD   = RGBColor(0xF5, 0xA6, 0x23)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xEA, 0xF0, 0xFB)

# ── Topic → search keywords ─────────────────────────────────────────────────
_KW = {
    "robot":        "industrial robot arm factory",
    "automotive":   "car manufacturing assembly line",
    "ev":           "electric vehicle charging",
    "electric":     "electric vehicle battery",
    "cnc":          "cnc machine manufacturing",
    "welding":      "welding sparks metal",
    "solar":        "solar panels renewable",
    "iot":          "smart factory sensors",
    "drone":        "drone technology",
    "3d print":     "3d printing prototype",
    "ai":           "artificial intelligence",
    "python":       "programming code",
    "data":         "data analytics charts",
    "mechatronics": "automation engineering",
    "plc":          "industrial control panel",
    "hydraulic":    "machinery engineering",
    "electronics":  "circuit board",
    "safety":       "industrial safety",
    "training":     "vocational workshop",
    "career":       "professional success",
    "programming":  "software development",
    "networking":   "network server",
    "mechanical":   "mechanical gears",
    "electrical":   "electrical panel",
    "manufacturing": "factory production",
    "battery":      "battery technology",
    "renewable":    "renewable energy",
    "healthcare":   "healthcare medical hospital",
    "machine learning": "machine learning ai technology",
    "ml":           "machine learning ai technology",
}

# Pexels API key (optional - get from https://www.pexels.com/api/)
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY", "")

def _keywords(topic: str) -> str:
    tl = topic.lower()
    for key, kw in _KW.items():
        if key in tl:
            return kw
    return "technology industry"

def _fetch_lorem_picsum(topic: str, slide_idx: int = 0, width: int = 800, height: int = 600):
    """Fetch unique image for each slide using different seeds"""
    try:
        # Create unique seed combining topic and slide index
        unique_string = f"{topic}_slide_{slide_idx}_{random.randint(1, 10000)}"
        seed = hashlib.md5(unique_string.encode()).hexdigest()[:10]
        
        url = f"https://picsum.photos/seed/{seed}/{width}/{height}"
        path = os.path.join(CACHE_DIR, f"picsum_{seed}.jpg")
        
        if os.path.exists(path) and os.path.getsize(path) > 8000:
            return path
            
        req = urllib.request.Request(url, headers={
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        })
        with urllib.request.urlopen(req, timeout=10) as r:
            data = r.read()
            if len(data) > 8000:
                with open(path, "wb") as f:
                    f.write(data)
                return path
    except Exception as e:
        print(f"Picsum fetch error: {e}")
    return None

def _fetch_pexels(topic: str, slide_idx: int = 0, width: int = 800, height: int = 600):
    """Try Pexels API with pagination for different images per slide"""
    if not PEXELS_API_KEY:
        return None
        
    try:
        kw = _keywords(topic)
        page = (slide_idx % 10) + 1
        url = f"https://api.pexels.com/v1/search?query={urllib.parse.quote(kw)}&per_page=1&page={page}&orientation=landscape"
        
        cache_key = hashlib.md5(f"{url}_{slide_idx}".encode()).hexdigest()
        path = os.path.join(CACHE_DIR, f"pexels_{cache_key}.jpg")
        
        if os.path.exists(path) and os.path.getsize(path) > 8000:
            return path
        
        req = urllib.request.Request(url, headers={
            "Authorization": PEXELS_API_KEY,
            "User-Agent": "Mozilla/5.0"
        })
        with urllib.request.urlopen(req, timeout=10) as r:
            data = json.loads(r.read().decode())
            photos = data.get('photos', [])
            if photos:
                img_url = photos[0]['src']['medium']
                img_req = urllib.request.Request(img_url, headers={"User-Agent": "Mozilla/5.0"})
                with urllib.request.urlopen(img_req, timeout=10) as img_r:
                    img_data = img_r.read()
                    if len(img_data) > 8000:
                        with open(path, "wb") as f:
                            f.write(img_data)
                        return path
    except Exception as e:
        print(f"Pexels fetch error: {e}")
    return None

def _fetch_image(topic: str, slide_idx: int = 0, width: int = 800, height: int = 600):
    """Fetch image with slide-specific uniqueness"""
    if PEXELS_API_KEY:
        result = _fetch_pexels(topic, slide_idx, width, height)
        if result:
            return result
    return _fetch_lorem_picsum(topic, slide_idx, width, height)

# ── Icons and Pills ─────────────────────────────────────────────────────────
_ICONS = {
    "robot": "🤖", "automotive": "🚗", "ev": "⚡", "electric": "⚡",
    "cnc": "⚙️", "weld": "🔧", "solar": "☀️", "iot": "📡",
    "drone": "🛸", "3d": "🖨️", "ai": "🧠", "data": "📊",
    "plc": "🏭", "safety": "🦺", "electronics": "💡", "career": "🏆",
    "mechatronics": "⚙️", "hydraulic": "🔩", "pneumatic": "💨",
    "programming": "💻", "networking": "🌐", "mechanical": "🔧",
    "electrical": "⚡", "manufacturing": "🏭", "battery": "🔋",
    "renewable": "🌱", "training": "🎓", "python": "🐍",
    "healthcare": "🏥", "medical": "🏥", "hospital": "🏥",
    "machine learning": "🤖", "ml": "🤖",
}

_PILLS = {
    "robot":       ["ABB · FANUC · KUKA",      "ISO 10218 Safety Std",    "90% Placement Rate"],
    "automotive":  ["Maruti · Tata · Bosch",    "₹25K–45K Starting",       "Top Industry Demand"],
    "ev":          ["BMS · BLDC · CAN Bus",     "ARAI Certification",       "₹28K Avg Salary"],
    "cnc":         ["Fanuc · Siemens 840D",     "G-Code & M-Code",          "₹25K Avg Salary"],
    "weld":        ["MIG · TIG · Plasma",       "AWS D1.1 Standard",        "₹22K Avg Salary"],
    "plc":         ["Siemens · Allen-Bradley",  "IEC 61131-3 Standard",     "₹30K Avg Salary"],
    "solar":       ["NABCEP Certified",         "On-Grid · Off-Grid",       "₹20K Avg Salary"],
    "iot":         ["MQTT · OPC-UA · REST",     "Industry 4.0 Ready",       "₹35K Avg Salary"],
    "data":        ["Python · SQL · Power BI",  "NASSCOM Recognised",       "₹40K Avg Salary"],
    "programming": ["Python · Java · C++",      "Full Stack Ready",         "₹35K Avg Salary"],
    "mechanical":  ["CAD · CAM · CAE",          "AutoCAD Certified",        "₹25K Avg Salary"],
    "electrical":  ["Panel Design · Wiring",    "BEE Certified",            "₹22K Avg Salary"],
    "healthcare":  ["NABH Certified",           "Nursing Council Approved", "₹40K Avg Salary"],
    "machine learning": ["TensorFlow · PyTorch", "AI/ML Certified",         "₹50K Avg Salary"],
}

def _icon_for(topic: str) -> str:
    tl = topic.lower()
    for k, v in _ICONS.items():
        if k in tl:
            return v
    return "🎓"

def _pills_for(topic: str):
    tl = topic.lower()
    for k, v in _PILLS.items():
        if k in tl:
            return v
    return ["Industry 4.0 Ready", "SkillSonics Certified", "Job Placement Support"]

# ── Draw branded fallback panel ─────────────────────────────────────────────
def _draw_panel(slide, topic: str, left: float, top: float, w: float, h: float, slide_idx: int = 0):
    """Rich illustrated panel when no internet photo available"""
    # Main background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(w), Inches(h),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_NAVY
    bg.line.fill.background()

    # Decorative circles with variation per slide
    colors = [C_GOLD, C_GREEN, RGBColor(0xFF, 0xFF, 0xFF)]
    positions = [
        (left + w - 1.35, top - 0.35, 1.7, 1.7),
        (left - 0.25, top + h - 1.1, 1.3, 1.3),
        (left + w * 0.3, top + h * 0.55, 0.5, 0.5),
    ]
    
    for i, (px, py, pw, ph) in enumerate(positions):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px), Inches(py), Inches(pw), Inches(ph))
        c.fill.solid()
        c.fill.fore_color.rgb = colors[i % len(colors)]
        c.line.fill.background()

    # Big emoji icon
    icon = _icon_for(topic)
    itb = slide.shapes.add_textbox(Inches(left), Inches(top + 0.35), Inches(w), Inches(1.3))
    itf = itb.text_frame
    itf.text = icon
    ip = itf.paragraphs[0]
    ip.font.size = Pt(64)
    ip.alignment = PP_ALIGN.CENTER

    # Topic label
    short = (topic[:30] + "…") if len(topic) > 30 else topic
    ltb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 1.75), Inches(w - 0.5), Inches(0.85))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    ltf.text = short
    lp = ltf.paragraphs[0]
    lp.font.size = Pt(18)
    lp.font.bold = True
    lp.font.color.rgb = C_WHITE
    lp.alignment = PP_ALIGN.CENTER

    # Gold divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + 0.5), Inches(top + 2.72),
        Inches(w - 1.0), Inches(0.04),
    )
    div.fill.solid()
    div.fill.fore_color.rgb = C_GOLD
    div.line.fill.background()

    # Sub-label
    stb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 2.82), Inches(w - 0.5), Inches(0.42))
    stf = stb.text_frame
    stf.text = "SkillSonics Training Program"
    sp = stf.paragraphs[0]
    sp.font.size = Pt(11)
    sp.font.color.rgb = C_GOLD
    sp.alignment = PP_ALIGN.CENTER

    # Stat pills
    pills = _pills_for(topic)
    for j, pill in enumerate(pills[:3]):
        py = top + 3.38 + j * 0.62
        pb = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.35), Inches(py),
            Inches(w - 0.7), Inches(0.48),
        )
        pb.fill.solid()
        pb.fill.fore_color.rgb = C_WHITE
        pb.line.fill.background()

        ptb = slide.shapes.add_textbox(
            Inches(left + 0.4), Inches(py + 0.07),
            Inches(w - 0.8), Inches(0.36),
        )
        ptf = ptb.text_frame
        ptf.text = pill
        pp_ = ptf.paragraphs[0]
        pp_.font.size = Pt(11)
        pp_.font.bold = True
        pp_.font.color.rgb = C_NAVY
        pp_.alignment = PP_ALIGN.CENTER

# ══════════════════════════════════════════════════════════════════════════════
class VisualEnhancer:

    def get_image_for_topic(self, topic: str, slide_idx: int = 0):
        return _fetch_image(topic, slide_idx)

    def add_image(self, slide, image_path, left, top, width, height):
        if not image_path or not os.path.exists(image_path):
            return False
        try:
            slide.shapes.add_picture(
                image_path,
                Inches(left), Inches(top), Inches(width), Inches(height),
            )
            return True
        except Exception as e:
            print(f"Add image error: {e}")
            return False

    def add_visual_right(self, slide, topic: str, slide_idx: int = 0, top: float = 1.3, height: float = 5.6) -> bool:
        """Place unique photo or branded panel on the RIGHT half"""
        img = _fetch_image(topic, slide_idx)
        if img and self.add_image(slide, img, left=5.1, top=top, width=4.65, height=height):
            return True
        _draw_panel(slide, topic, left=5.1, top=top, w=4.65, h=height, slide_idx=slide_idx)
        return True

    def add_visual_left(self, slide, topic: str, slide_idx: int = 0, top: float = 1.3, height: float = 5.6) -> bool:
        """Place unique photo or branded panel on the LEFT half"""
        img = _fetch_image(topic, slide_idx)
        if img and self.add_image(slide, img, left=0.25, top=top, width=4.65, height=height):
            return True
        _draw_panel(slide, topic, left=0.25, top=top, w=4.65, h=height, slide_idx=slide_idx)
        return True

    def add_background_image(self, slide, topic: str, slide_idx: int = 0):
        """Add full-slide background image"""
        img = _fetch_image(topic, slide_idx, width=1280, height=960)
        if not img:
            return False
        try:
            pic = slide.shapes.add_picture(img, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
            spTree = slide.shapes._spTree
            sp = pic._element
            spTree.remove(sp)
            spTree.insert(2, sp)
            return True
        except Exception as e:
            print(f"Background image error: {e}")
            return False