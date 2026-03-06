"""
SkillSonics Slide Templates — Gamma-quality design with unique images per slide.
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from config import FONTS, ORG_INFO

SW, SH = 10.0, 7.5

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0F, 0x2A, 0x4A)
NAVY2   = RGBColor(0x1A, 0x3D, 0x6B)
NAVY3   = RGBColor(0x0A, 0x1E, 0x38)
GOLD    = RGBColor(0xF5, 0xA6, 0x23)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
TEAL    = RGBColor(0x06, 0xB6, 0xD4)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
RED     = RGBColor(0xEF, 0x44, 0x44)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFF     = RGBColor(0xF4, 0xF6, 0xFA)
SLATE   = RGBColor(0xF1, 0xF5, 0xF9)
DARK    = RGBColor(0x0D, 0x1B, 0x2A)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
BORDER  = RGBColor(0xE2, 0xE8, 0xF0)

# Card fill + accent pairs for content grids
CARD_BG  = [RGBColor(0xEF,0xF6,0xFF), RGBColor(0xFF,0xF9,0xEB),
            RGBColor(0xEC,0xFD,0xF5), RGBColor(0xF5,0xF0,0xFF)]
CARD_ACC = [NAVY, GOLD, GREEN, PURPLE]
STAT_ACC = [NAVY, GOLD, GREEN, PURPLE]

# Compat exports expected by old ppt_builder imports
HEADER_H   = 1.12
BRAND_BAR_Y= 7.04
CONTENT_TOP= 1.22
CONTENT_H  = 5.72
MARGIN     = 0.34

def _rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _title_pt(text, cap=28):
    n = len(text)
    if n <= 30: return min(cap, 28)
    if n <= 50: return min(cap, 22)
    return min(cap, 18)

def _bullet_pt(n, has_visual=True):
    if n <= 3: return 18
    if n <= 4: return 16
    if n <= 5: return 14
    return 12

def _shadow(shape, opacity=10, blur=50000, dist=20000):
    try:
        sp   = shape._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is None: return
        old = spPr.find(qn("a:effectLst"))
        if old is not None: spPr.remove(old)
        eff = etree.SubElement(spPr, qn("a:effectLst"))
        s   = etree.SubElement(eff, qn("a:outerShdw"))
        s.set("blurRad", str(blur)); s.set("dist", str(dist))
        s.set("dir", "5400000");     s.set("algn", "ctr")
        c = etree.SubElement(s, qn("a:srgbClr")); c.set("val", "000000")
        a = etree.SubElement(c, qn("a:alpha"));   a.set("val", str(opacity*1000))
    except Exception:
        pass

class SlideTemplate:
    def __init__(self, prs):
        self.prs   = prs
        self.fonts = FONTS
        self.colors = {"primary": "#0F2A4A", "secondary": "#F5A623", "accent": "#10B981"}

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    # ── Core primitives ───────────────────────────────────────────────────────
    def _rect(self, slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
        s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
        return s

    def _rr(self, slide, x, y, w, h, fill):
        return self._rect(slide, x, y, w, h, fill, MSO_SHAPE.ROUNDED_RECTANGLE)

    def _oval(self, slide, x, y, w, h, fill):
        return self._rect(slide, x, y, w, h, fill, MSO_SHAPE.OVAL)

    def _tb(self, slide, x, y, w, h, text, size, bold=False,
            color=None, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap; tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color or DARK
        p.font.name = "Calibri"; p.alignment = align
        return tb

    def _footer(self, slide):
        self._rect(slide, 0, 7.04, SW, 0.46, DARK)
        self._tb(slide, 0.4, 7.1, 9.2, 0.34,
                 f"{ORG_INFO['name']}   ·   {ORG_INFO['tagline']}   ·   {ORG_INFO['website']}",
                 9, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.LEFT)

    def _page_header(self, slide, title, bg=NAVY):
        """Standard top bar used on all non-title slides."""
        self._rect(slide, 0, 0, SW, SH, OFF)        # page bg
        self._rect(slide, 0, 0, SW, HEADER_H, bg)   # header bar
        self._rect(slide, 0, 0, 0.28, HEADER_H, GOLD)  # gold left accent
        pt = _title_pt(title, 28)
        self._tb(slide, 0.52, 0.15, 9.0, HEADER_H-0.18,
                 title, pt, bold=True, color=WHITE)

    # Compat: old builder calls these
    def add_brand_bar(self, slide):   self._footer(slide)
    def add_brand_shape(self, slide): self._footer(slide)

    # ══════════════════════════════════════════════════════════════════════════
    # TITLE SLIDE — with optional background image
    # ══════════════════════════════════════════════════════════════════════════
    def create_title_slide(self, title, subtitle="", image_path=None, topic="", visual=None, slide_idx=0):
        slide = self._blank()

        # Try to add background image if visual and topic provided
        if visual and topic:
            try:
                visual.add_background_image(slide, topic, slide_idx)
            except:
                pass
        
        # Dark navy base (overlay if image added, or solid if not)
        base = self._rect(slide, 0, 0, SW, SH, NAVY)
        if visual and topic:
            # Make semi-transparent overlay
            try:
                base.fill.solid()
                base.fill.fore_color.rgb = NAVY
            except:
                pass

        # Layered geometric circles for depth
        self._oval(slide,  6.6, -1.2, 5.6, 5.6, NAVY2)
        self._oval(slide,  7.8,  3.8, 4.2, 4.2, NAVY3)
        self._oval(slide, -1.0,  5.2, 3.8, 3.8, NAVY3)
        self._oval(slide,  0.4, -0.6, 2.2, 2.2, GOLD)      # gold top-left accent
        self._oval(slide,  8.8,  6.2, 1.0, 1.0, GREEN)     # green bottom-right dot

        # Thin gold top stripe
        self._rect(slide, 0, 0, SW, 0.10, GOLD)

        # SkillSonics badge
        self._rr(slide, 0.44, 0.26, 2.6, 0.48, NAVY2)
        self._tb(slide, 0.54, 0.26, 2.5, 0.48,
                 f"🎓  {ORG_INFO['name']}", 11, bold=True, color=GOLD)

        # Main title
        tpt = 52 if len(title) <= 20 else (42 if len(title) <= 34 else 32)
        self._tb(slide, 0.44, 1.1, 9.0, 3.0,
                 title, tpt, bold=True, color=WHITE, wrap=True)

        # Gold rule
        self._rect(slide, 0.44, 4.35, 2.0, 0.08, GOLD)

        # Subtitle
        if subtitle:
            spt = 20 if len(subtitle) <= 65 else 16
            self._tb(slide, 0.44, 4.52, 8.4, 1.4,
                     subtitle, spt, color=RGBColor(0xB0,0xC4,0xDE), wrap=True)

        # Bottom info bar
        self._rect(slide, 0, 6.6, SW, 0.9, NAVY3)
        self._tb(slide, 0.44, 6.70, 9.0, 0.46,
                 f"{ORG_INFO['tagline']}   ·   {ORG_INFO['website']}   ·   {ORG_INFO['phone']}",
                 9.5, color=RGBColor(0x7A,0x92,0xB0))

        return slide

    # ══════════════════════════════════════════════════════════════════════════
    # CONTENT SLIDE — 2-column card grid OR single column with right image
    # ══════════════════════════════════════════════════════════════════════════
    def create_content_slide(self, title, points, image_path=None,
                              visual=None, topic=None, slide_idx=0):
        slide = self._blank()
        self._page_header(slide, title)

        n = min(len(points), 6)
        if n == 0:
            self._footer(slide)
            return slide

        # Check if we should add a visual on the right
        has_visual = visual is not None and topic is not None
        
        if has_visual:
            # Add visual on the right half with unique image per slide
            visual.add_visual_right(slide, topic, slide_idx=slide_idx, top=1.3, height=5.6)
            
            # Single column layout on left side (narrower)
            cols = 1
            rows = n
            M = MARGIN
            gap = 0.20
            ct = CONTENT_TOP + 0.06
            avail_w = 4.8  # Left side width only
            avail_h = 7.04 - ct - 0.10
            cw = avail_w - 0.2
            ch = (avail_h - gap*(rows-1)) / rows

            for i, point in enumerate(points[:6]):
                y = ct + i*(ch+gap)
                x = M
                acc = CARD_ACC[i % 4]
                bg = CARD_BG[i % 4]

                card = self._rr(slide, x, y, cw, ch, bg)
                _shadow(card, opacity=8, blur=35000, dist=12000)

                # Left coloured accent strip
                self._rect(slide, x, y, 0.07, ch, acc)

                # Numbered circle
                self._oval(slide, x+0.18, y+ch/2-0.26, 0.52, 0.52, acc)
                self._tb(slide, x+0.18, y+ch/2-0.22, 0.52, 0.44,
                         str(i+1), 15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

                # Bullet text
                fpt = 12 if len(point) > 95 else (13 if len(point) > 65 else 14)
                self._tb(slide, x+0.84, y+0.10,
                         cw-0.96, ch-0.20, point, fpt, color=DARK, wrap=True)
        else:
            # Original 2-column layout when no visual
            cols = 2
            rows = -(-n // cols)
            M = MARGIN
            gap = 0.20
            ct = CONTENT_TOP + 0.06
            avail_w = SW - 2*M
            avail_h = 7.04 - ct - 0.10
            cw = (avail_w - gap*(cols-1)) / cols
            ch = (avail_h - gap*(rows-1)) / rows

            for i, point in enumerate(points[:6]):
                row = i // cols; col = i % cols
                x = M + col*(cw+gap)
                y = ct + row*(ch+gap)
                acc = CARD_ACC[i % 4]
                bg = CARD_BG[i % 4]

                card = self._rr(slide, x, y, cw, ch, bg)
                _shadow(card, opacity=8, blur=35000, dist=12000)

                # Left coloured accent strip
                self._rect(slide, x, y, 0.07, ch, acc)

                # Numbered circle
                self._oval(slide, x+0.18, y+ch/2-0.26, 0.52, 0.52, acc)
                self._tb(slide, x+0.18, y+ch/2-0.22, 0.52, 0.44,
                         str(i+1), 15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

                # Bullet text
                fpt = 12 if len(point) > 95 else (13 if len(point) > 65 else 14)
                self._tb(slide, x+0.84, y+0.10,
                         cw-0.96, ch-0.20, point, fpt, color=DARK, wrap=True)

        self._footer(slide)
        return slide

    # ══════════════════════════════════════════════════════════════════════════
    # STATISTICS SLIDE — bold number cards, full height, coloured tops
    # ══════════════════════════════════════════════════════════════════════════
    def create_stats_slide(self, title, statistics):
        slide = self._blank()
        self._page_header(slide, title)

        n = min(len(statistics), 4)
        if n == 0:
            self._footer(slide)
            return slide

        M = MARGIN
        gap = 0.20
        cw = (SW - 2*M - gap*(n-1)) / n
        ct = CONTENT_TOP + 0.08
        ch = 7.04 - ct - 0.10

        STAT_BG = [RGBColor(0xEF,0xF6,0xFF), RGBColor(0xFF,0xF9,0xEB),
                   RGBColor(0xEC,0xFD,0xF5), RGBColor(0xF5,0xF0,0xFF)]

        for i, stat in enumerate(statistics[:4]):
            x = M + i*(cw+gap)
            col = STAT_ACC[i % 4]
            num = stat.get("number", "—")
            ctx = stat.get("context", "")

            card = self._rr(slide, x, ct, cw, ch, STAT_BG[i % 4])
            _shadow(card, opacity=10, blur=50000, dist=18000)

            # Thick coloured top bar
            self._rr(slide, x, ct, cw, 0.36, col)

            # Big number
            npt = 56 if len(num) <= 5 else (42 if len(num) <= 8 else 30)
            self._tb(slide, x+0.08, ct+0.44, cw-0.16, 1.55,
                     num, npt, bold=True, color=col, align=PP_ALIGN.CENTER, wrap=False)

            # Pill divider
            self._rr(slide, x+cw*0.25, ct+2.10, cw*0.50, 0.07, col)

            # Context
            cpt = 11 if len(ctx) > 70 else 13
            self._tb(slide, x+0.14, ct+2.26, cw-0.28, ch-2.36,
                     ctx, cpt, color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

        self._footer(slide)
        return slide

    # ══════════════════════════════════════════════════════════════════════════
    # COMPARISON SLIDE — two full-height panels (light vs dark)
    # ══════════════════════════════════════════════════════════════════════════
    def create_comparison_slide(self, title, left_title, left_content,
                                 right_title, right_content):
        slide = self._blank()
        self._page_header(slide, title)

        M = MARGIN
        ct = CONTENT_TOP + 0.06
        ch = 7.04 - ct - 0.10
        half = (SW - 2*M - 0.22) / 2

        # LEFT — white card
        lp = self._rr(slide, M, ct, half, ch, WHITE)
        _shadow(lp, opacity=10)
        self._rr(slide, M, ct, half, 0.55, SLATE)
        self._tb(slide, M+0.2, ct+0.10, half-0.3, 0.42,
                 f"❌  {left_title}", 13, bold=True, color=MUTED)

        bpt = 12 if len(left_content) > 4 else 13
        ih = (ch - 0.65) / max(len(left_content[:5]), 1)
        ty = ct + 0.65
        for item in left_content[:5]:
            self._oval(slide, M+0.18, ty+ih*0.32, 0.14, 0.14, RED)
            self._tb(slide, M+0.42, ty+0.04, half-0.58, ih-0.08,
                     item, bpt, color=DARK, wrap=True)
            ty += ih

        # RIGHT — navy card
        rx = M + half + 0.22
        rp = self._rr(slide, rx, ct, half, ch, NAVY)
        _shadow(rp, opacity=14)
        self._rr(slide, rx, ct, half, 0.55, NAVY3)
        self._tb(slide, rx+0.2, ct+0.10, half-0.3, 0.42,
                 f"✅  {right_title}", 13, bold=True, color=GOLD)

        ih = (ch - 0.65) / max(len(right_content[:5]), 1)
        ty = ct + 0.65
        for item in right_content[:5]:
            self._oval(slide, rx+0.18, ty+ih*0.32, 0.14, 0.14, GREEN)
            self._tb(slide, rx+0.42, ty+0.04, half-0.58, ih-0.08,
                     item, bpt, color=WHITE, wrap=True)
            ty += ih

        self._footer(slide)
        return slide

    # ══════════════════════════════════════════════════════════════════════════
    # PROCESS SLIDE — numbered cards (horizontal ≤4, vertical 5)
    # ══════════════════════════════════════════════════════════════════════════
    def create_process_slide(self, title, steps):
        slide = self._blank()
        self._page_header(slide, title)

        n = min(len(steps), 5)
        COLS = [NAVY, GOLD, GREEN, TEAL, PURPLE]
        M = MARGIN
        ct = CONTENT_TOP + 0.06
        ch = 7.04 - ct - 0.10
        gap = 0.18

        if n <= 4:
            cw = (SW - 2*M - gap*(n-1)) / n
            bh = ch - 0.75

            for i, step in enumerate(steps[:n]):
                x = M + i*(cw+gap)
                col = COLS[i % 5]

                # Number circle above card
                self._oval(slide, x+cw/2-0.34, ct+0.04, 0.68, 0.68, col)
                self._tb(slide, x+cw/2-0.34, ct+0.10, 0.68, 0.52,
                         str(i+1), 22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

                # Arrow connector
                if i < n-1:
                    ax = x+cw+0.02; ay = ct+0.34
                    self._rect(slide, ax, ay, gap-0.02, 0.06, BORDER)

                # Step card
                cy = ct+0.82
                card = self._rr(slide, x, cy, cw, bh, WHITE)
                _shadow(card, opacity=10)
                self._rr(slide, x, cy, cw, 0.24, col)

                fpt = 11 if len(step) > 90 else (12 if len(step) > 60 else 13)
                self._tb(slide, x+0.10, cy+0.30, cw-0.20, bh-0.38,
                         step, fpt, color=DARK, align=PP_ALIGN.CENTER, wrap=True)

        else:  # vertical layout for 5 steps
            sh_ = (ch - gap*(n-1)) / n
            for i, step in enumerate(steps[:n]):
                y = ct + i*(sh_+gap)
                col = COLS[i % 5]

                # Number pill
                self._oval(slide, M, y+sh_/2-0.28, 0.56, 0.56, col)
                self._tb(slide, M, y+sh_/2-0.24, 0.56, 0.46,
                         str(i+1), 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

                # Connector line
                if i < n-1:
                    self._rect(slide, M+0.26, y+sh_/2+0.28,
                               0.04, gap+sh_/2-0.30, BORDER)

                # Step card
                card = self._rr(slide, M+0.72, y+0.06,
                                SW-M-0.72-M, sh_-0.12, WHITE)
                _shadow(card, opacity=10)
                self._rect(slide, M+0.72, y+0.06, 0.08, sh_-0.12, col)

                fpt = 11 if len(step) > 90 else (13 if len(step) > 55 else 14)
                self._tb(slide, M+0.94, y+0.14,
                         SW-M-0.94-M, sh_-0.28, step, fpt, color=DARK, wrap=True)

        self._footer(slide)
        return slide

    # ══════════════════════════════════════════════════════════════════════════
    # TWO-COLUMN SLIDE
    # ══════════════════════════════════════════════════════════════════════════
    def create_two_column_slide(self, title, left_content, right_content,
                                 left_title="", right_title=""):
        slide = self._blank()
        self._page_header(slide, title)

        M = MARGIN; ct = CONTENT_TOP+0.06; ch = 7.04-ct-0.10
        cw = (SW - 2*M - 0.22) / 2

        for xi, (x, items, ttl, col) in enumerate([
            (M,            left_content,  left_title,  NAVY),
            (M+cw+0.22,    right_content, right_title, GREEN),
        ]):
            card = self._rr(slide, x, ct, cw, ch, WHITE)
            _shadow(card, opacity=10)
            self._rr(slide, x, ct, cw, 0.30, col)

            if ttl:
                self._tb(slide, x+0.2, ct+0.36, cw-0.35, 0.44,
                         ttl, 14, bold=True, color=col)

            bpt = 12 if len(items) > 4 else 14
            ty = ct + (0.86 if ttl else 0.42)
            ih = (ch - (0.9 if ttl else 0.45)) / max(len(items[:6]), 1)
            for item in items[:6]:
                self._oval(slide, x+0.18, ty+ih*0.30, 0.14, 0.14, col)
                self._tb(slide, x+0.42, ty+0.04, cw-0.56, ih-0.08,
                         item, bpt, color=DARK, wrap=True)
                ty += ih

        self._footer(slide)
        return slide

    # ══════════════════════════════════════════════════════════════════════════
    # SECTION HEADER SLIDE — with optional background image
    # ══════════════════════════════════════════════════════════════════════════
    def create_section_slide(self, title, subtitle="", visual=None, topic=None, slide_idx=0):
        slide = self._blank()
        
        # Try background image first
        if visual and topic:
            try:
                visual.add_background_image(slide, topic, slide_idx)
            except:
                pass
        
        # Base background
        self._rect(slide, 0, 0, SW, SH, NAVY)
        
        self._oval(slide, 7.0, -0.8, 4.5, 4.5, NAVY2)
        self._oval(slide, -0.6, 5.0, 3.0, 3.0, NAVY3)
        self._rect(slide, 0, 0, SW, 0.10, GOLD)
        
        tpt = 44 if len(title) <= 24 else (34 if len(title) <= 40 else 26)
        self._tb(slide, 0.6, 2.4, 8.8, 2.0,
                 title, tpt, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        
        if subtitle:
            self._tb(slide, 0.6, 4.55, 8.0, 1.0,
                     subtitle, 18, color=RGBColor(0xB0,0xC4,0xDE))
        
        self._footer(slide)
        return slide

    # ══════════════════════════════════════════════════════════════════════════
    # END / CONTACT SLIDE
    # ══════════════════════════════════════════════════════════════════════════
    def create_end_slide(self, title="Start Your Journey", contact_info=True):
        slide = self._blank()
        self._rect(slide, 0, 0, SW, SH, NAVY)
        self._oval(slide, 7.0, -1.0, 5.2, 5.2, NAVY2)
        self._oval(slide, -0.8, 5.2, 3.2, 3.2, RGBColor(0x06,0x4E,0x3B))
        self._oval(slide, 0.6, -0.5, 2.0, 2.0, GOLD)
        self._oval(slide, 8.8, 6.3, 1.0, 1.0, TEAL)
        self._rect(slide, 0, 0, SW, 0.10, GOLD)

        tpt = 48 if len(title) <= 20 else (36 if len(title) <= 32 else 26)
        self._tb(slide, 1.0, 1.5, 8.0, 2.0,
                 title, tpt, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        self._rect(slide, 3.8, 3.65, 2.4, 0.08, GOLD)

        if contact_info:
            card = self._rr(slide, 1.6, 3.88, 6.8, 2.7, WHITE)
            _shadow(card, opacity=16)
            for i, (icon, text) in enumerate([
                ("📧", ORG_INFO["email"]),
                ("🌐", ORG_INFO["website"]),
                ("📞", ORG_INFO["phone"]),
            ]):
                self._tb(slide, 2.1, 4.08+i*0.78, 6.0, 0.70,
                         f"{icon}   {text}", 16, color=DARK, align=PP_ALIGN.CENTER)

        return slide